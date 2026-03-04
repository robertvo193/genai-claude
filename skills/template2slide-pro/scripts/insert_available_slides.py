#!/usr/bin/env python3
"""
Insert Available slides into generated presentation

Usage:
    python insert_available_slides.py <generated.pptx> <AvailableSlide11.pptx> <output.pptx>

Example:
    python insert_available_slides.py ../../output_bromma/presentation.pptx AvailableSlide11.pptx ../../output_bromma/presentation_complete.pptx
"""

import sys
from pptx import Presentation

def insert_available_slides(generated_path, available_path, output_path):
    """
    Insert Available slides into generated presentation

    Structure:
    1. Slide 1 from generated (Title Slide)
    2. Available Slides 2-10 (9 slides)
    3. Remaining slides from generated (Slides 2 onwards)
    4. Available Slides 11-25 (15 slides)
    """
    print("📊 Inserting Available slides into generated presentation...\n")

    # Load presentations
    print(f"Loading generated presentation: {generated_path}")
    generated_ppt = Presentation(generated_path)
    print(f"  ✓ Loaded {len(generated_ppt.slides)} slides")

    print(f"Loading Available slides: {available_path}")
    available_ppt = Presentation(available_path)
    print(f"  ✓ Loaded {len(available_ppt.slides)} slides")

    # Create output presentation
    output_ppt = Presentation()

    # Step 1: Copy Slide 1 (Title Slide) from generated
    print("\n📌 Step 1: Copying Slide 1 (Title Slide) from generated presentation")
    _copy_slide(generated_ppt.slides[0], output_ppt.slides.add_slide())
    print("  ✓ Slide 1 copied")

    # Step 2: Insert Available Slides 2-10 (indices 1-9 in 0-based)
    print("\n📌 Step 2: Inserting Available Slides 2-10 after Slide 1")
    available_count = min(9, len(available_ppt.slides) - 1)
    for i in range(available_count):
        source_idx = i + 1  # Skip first available slide (index 0)
        target_slide = output_ppt.slides.add_slide()
        _copy_slide(available_ppt.slides[source_idx], target_slide)
        print(f"  ✓ Available Slide {source_idx + 1} inserted (Position {len(output_ppt.slides)})")

    # Step 3: Copy remaining slides from generated (Slides 2 onwards, index 1+)
    print("\n📌 Step 3: Copying remaining slides from generated presentation")
    for i in range(1, len(generated_ppt.slides)):
        target_slide = output_ppt.slides.add_slide()
        _copy_slide(generated_ppt.slides[i], target_slide)
        print(f"  ✓ Generated Slide {i + 1} copied (Position {len(output_ppt.slides)})")

    # Step 4: Insert Available Slides 11-25 (indices 10-24 in 0-based)
    print("\n📌 Step 4: Inserting Available Slides 11-25 at the end")
    for i in range(10, min(25, len(available_ppt.slides))):
        if i < len(available_ppt.slides):
            target_slide = output_ppt.slides.add_slide()
            _copy_slide(available_ppt.slides[i], target_slide)
            print(f"  ✓ Available Slide {i + 1} inserted (Position {len(output_ppt.slides)})")

    # Save output presentation
    print(f"\n💾 Saving combined presentation to: {output_path}")
    output_ppt.save(output_path)

    print(f"\n✅ Successfully created combined presentation!")
    print(f"   Total slides: {len(output_ppt.slides)}")
    print(f"   - Slide 1: Generated title slide")
    print(f"   - Slides 2-{available_count + 1}: Available slides (2-10)")
    print(f"   - Slides {available_count + 2}-{available_count + len(generated_ppt.slides)}: Generated content slides")
    print(f"   - Slides {available_count + len(generated_ppt.slides) + 1}-{len(output_ppt.slides)}: Available slides (11-25)")


def _copy_slide(source_slide, target_slide):
    """
    Copy content from source slide to target slide

    Note: python-pptx doesn't have built-in slide cloning, so we manually copy shapes
    This copies shapes, text, images, and maintains the layout
    """
    # Clear default layouts from target slide
    for shape in target_slide.shapes:
        sp = shape.element
        sp.getparent().remove(sp)

    # Copy each shape from source to target
    for shape in source_slide.shapes:
        # Get the XML element from source shape
        shape_element = shape.element

        # Import it into target slide
        target_slide.shapes._spTree.insert_element_before(
            target_slide.shapes._spTree.ppr_lstChild,  # Insert before last child
            shape_element
        )


def main():
    if len(sys.argv) < 4:
        print("Usage: python insert_available_slides.py <generated.pptx> <AvailableSlide11.pptx> <output.pptx>")
        print("")
        print("Example:")
        print("  python insert_available_slides.py ../../output_bromma/presentation.pptx AvailableSlide11.pptx ../../output_bromma/presentation_complete.pptx")
        sys.exit(1)

    generated_path = sys.argv[1]
    available_path = sys.argv[2]
    output_path = sys.argv[3]

    try:
        insert_available_slides(generated_path, available_path, output_path)
        sys.exit(0)
    except Exception as e:
        print(f"\n❌ Error: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()
