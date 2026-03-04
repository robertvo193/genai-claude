#!/usr/bin/env python3
"""
Insert reference slides from System_architecture.pptx and Available_Slide.pptx into generated presentation.

Usage:
    python insert_reference_slides.py <generated_pptx> <project_info_json> [output_pptx]

This script:
1. Inserts architecture template slide from System_architecture.pptx after the architecture diagram slide
2. Inserts Available slides (2-10) after slide 1
3. Inserts Available slides (11-25) after the last slide
4. Copies background from generated slides to inserted slides
"""

import sys
import json
import shutil
from pathlib import Path
from copy import deepcopy
from lxml import etree

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    print("Error: python-pptx not installed. Install with: pip install python-pptx")
    sys.exit(1)


# Mapping deployment method to slide index in System_architecture.pptx
DEPLOYMENT_SLIDE_MAP = {
    'cloud': 1,  # Slide 2 (0-indexed: 1)
    'on-premise': 2,  # Slide 3 (0-indexed: 2)
    'on-prem': 2,  # Alias for on-premise
    'hybrid': 3,  # Slide 4 - Hybrid (AI Interference at site, Dashboard + Training Cloud)
    'hybrid-training-on-prem': 4,  # Slide 5 - Hybrid (AI Interference + Training at site, Dashboard Cloud)
    'hybrid-training-onprem': 4,  # Alias
    '4g-vpn-bridge': None,  # Not in template
    'vimov': None,  # Not in template
}


def get_deployment_method(project_info_path):
    """Extract deployment method from project_info JSON file"""
    try:
        with open(project_info_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        
        # Handle both nested and flat structures
        if 'project_info' in data:
            project_info = data['project_info']
        else:
            project_info = data
        
        deployment_method = project_info.get('deployment_method', '').lower()
        # Normalize deployment method names
        if 'cloud' in deployment_method:
            return 'cloud'
        elif 'hybrid' in deployment_method:
            # Check if training is on-premise
            if 'training' in deployment_method and ('on-prem' in deployment_method or 'onprem' in deployment_method):
                return 'hybrid-training-on-prem'
            else:
                return 'hybrid'
        elif 'on-prem' in deployment_method or 'onprem' in deployment_method:
            return 'on-premise'
        elif '4g' in deployment_method or 'vpn' in deployment_method:
            return '4g-vpn-bridge'
        elif 'vimov' in deployment_method:
            return 'vimov'
        
        return deployment_method
    except Exception as e:
        print(f"Warning: Could not read deployment method from {project_info_path}: {e}")
        return None


def find_architecture_slide_index(presentation):
    """Find the index of the architecture diagram slide (type: diagram)"""
    # The architecture slide is typically the one with "System Architecture" or "Proposed System Architecture" title
    # We'll look for slides with diagram-related content
    for i, slide in enumerate(presentation.slides):
        # Check slide notes or title for architecture indicators
        # Since we can't easily check slide type, we'll use a heuristic:
        # Look for slides with "Architecture" in title or notes
        try:
            # Check if slide has a title shape
            for shape in slide.shapes:
                if hasattr(shape, 'text') and 'architecture' in shape.text.lower():
                    return i
        except:
            pass
    
    # Default: assume architecture slide is around slide 4 (0-indexed: 3)
    # This is based on SLIDE_TEMPLATE.md which says System Architecture is Slide 4
    return 3


def copy_slide_background(source_slide, target_slide):
    """Copy background from source slide to target slide at XML level"""
    try:
        # Access the slide's XML element directly (avoid python-pptx API)
        source_spTree = source_slide._element
        target_spTree = target_slide._element

        # Find or create background element in target slide
        nsmap = {
            'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'
        }

        # Get background from source slide
        source_bg_el = source_spTree.find('.//p:bg', nsmap)

        if source_bg_el is not None:
            # Remove existing background from target slide
            target_bg_el = target_spTree.find('.//p:bg', nsmap)
            if target_bg_el is not None:
                target_bg_el.getparent().remove(target_bg_el)

            # Clone and insert source background
            new_bg_el = etree.fromstring(etree.tostring(source_bg_el))

            # Find the correct insertion point (before spTree)
            target_sld = target_spTree
            # Insert background as first child after clrMapOvr
            clr_map = target_sld.find('.//p:clrMapOvr', nsmap)
            if clr_map is not None:
                clr_map.getparent().insert(clr_map.getparent().index(clr_map) + 1, new_bg_el)
            else:
                # If no clrMapOvr, insert as first child
                target_sld.insert(0, new_bg_el)

        return True
    except Exception as e:
        print(f"Warning: Could not copy background: {e}")
        return False


def duplicate_slide(pres, source_index):
    """
    Duplicate a slide in the presentation (based on rearrange.py).
    Returns the new slide.
    """
    import six

    source = pres.slides[source_index]

    # Use source's layout to preserve formatting
    new_slide = pres.slides.add_slide(source.slide_layout)

    # Copy background element BEFORE clearing shapes (to preserve it)
    try:
        nsmap = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
        # Access the slide's element through the slide's part
        source_sld_el = source._element  # Use _element for direct XML access
        source_bg_el = source_sld_el.find('.//p:bg', nsmap)
        if source_bg_el is not None:
            # Clone background element by serializing to string and parsing back
            # This avoids deepcopy issues with python-pptx objects
            bg_xml_str = etree.tostring(source_bg_el, encoding='unicode')
            new_bg_el = etree.fromstring(bg_xml_str)
            # Insert into new slide
            new_sld_el = new_slide._element
            # Find correct insertion point (after clrMapOvr if exists)
            clr_map = new_sld_el.find('.//p:clrMapOvr', nsmap)
            if clr_map is not None:
                clr_map.getparent().insert(clr_map.getparent().index(clr_map) + 1, new_bg_el)
            else:
                # Insert as first child
                new_sld_el.insert(0, new_bg_el)
    except Exception as e:
        import traceback
        print(f"Warning: Could not copy background element: {e}")
        print(f"Traceback: {traceback.format_exc()}")
    
    # Collect all image and media relationships from the source slide
    image_rels = {}
    for rel_id, rel in six.iteritems(source.part.rels):
        if "image" in rel.reltype or "media" in rel.reltype:
            image_rels[rel_id] = rel
    
    # CRITICAL: Clear placeholder shapes to avoid duplicates
    for shape in list(new_slide.shapes):
        try:
            sp = shape.element
            sp.getparent().remove(sp)
        except:
            pass
    
    # Copy all shapes from source
    for shape in source.shapes:
        try:
            el = shape.element
            new_el = deepcopy(el)
            new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")
            
            # Handle picture shapes - need to update the blip reference
            blips = new_el.xpath(".//a:blip[@r:embed]")
            for blip in blips:
                old_rId = blip.get(
                    "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
                )
                if old_rId in image_rels:
                    # Create a new relationship in the destination slide for this image
                    old_rel = image_rels[old_rId]
                    new_rId = new_slide.part.rels.get_or_add(
                        old_rel.reltype, old_rel._target
                    )
                    # Update the blip's embed reference to use the new relationship ID
                    blip.set(
                        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed",
                        new_rId,
                    )
        except Exception as e:
            print(f"Warning: Could not copy shape: {e}")
    
    # Note: Background is preserved from source slide during shape copying above
    # No need for additional background copying

    return new_slide


def copy_slide_from_other_pres(source_pres, source_index, target_pres):
    """
    Copy a slide from source presentation to target presentation.
    Returns the new slide.
    """
    import six

    source_slide = source_pres.slides[source_index]

    # Use source slide's layout
    new_slide = target_pres.slides.add_slide(source_slide.slide_layout)

    # Copy background element BEFORE clearing shapes (to preserve it)
    try:
        nsmap = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
        # Access the slide's element through the slide's part
        source_sld_el = source_slide._element  # Use _element for direct XML access
        source_bg_el = source_sld_el.find('.//p:bg', nsmap)
        if source_bg_el is not None:
            # Clone background element by serializing to string and parsing back
            # This avoids deepcopy issues with python-pptx objects
            bg_xml_str = etree.tostring(source_bg_el, encoding='unicode')
            new_bg_el = etree.fromstring(bg_xml_str)
            # Insert into new slide
            new_sld_el = new_slide._element
            # Find correct insertion point (after clrMapOvr if exists)
            clr_map = new_sld_el.find('.//p:clrMapOvr', nsmap)
            if clr_map is not None:
                clr_map.getparent().insert(clr_map.getparent().index(clr_map) + 1, new_bg_el)
            else:
                # Insert as first child
                new_sld_el.insert(0, new_bg_el)
    except Exception as e:
        print(f"Warning: Could not copy background element: {e}")
    
    # Collect all image and media relationships from the source slide
    image_rels = {}
    for rel_id, rel in six.iteritems(source_slide.part.rels):
        if "image" in rel.reltype or "media" in rel.reltype:
            image_rels[rel_id] = rel
    
    # Clear existing shapes (placeholders)
    for shape in list(new_slide.shapes):
        try:
            sp = shape.element
            sp.getparent().remove(sp)
        except:
            pass
    
    # Copy all shapes from source slide
    for shape in source_slide.shapes:
        try:
            el = shape.element
            new_el = deepcopy(el)
            new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")
            
            # Handle picture shapes - update blip references
            blips = new_el.xpath(".//a:blip[@r:embed]")
            for blip in blips:
                old_rId = blip.get(
                    "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
                )
                if old_rId in image_rels:
                    old_rel = image_rels[old_rId]
                    new_rId = new_slide.part.rels.get_or_add(old_rel.reltype, old_rel._target)
                    blip.set(
                        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed",
                        new_rId,
                    )
        except Exception as e:
            print(f"Warning: Could not copy shape: {e}")
    
    # Note: Background is preserved from source slide during shape copying above
    # No need for additional background copying

    return new_slide


def insert_reference_slides(generated_pptx_path, project_info_path, output_pptx_path=None):
    """
    Insert reference slides into generated presentation.
    
    Args:
        generated_pptx_path: Path to the generated PowerPoint file
        project_info_path: Path to project_info.json file
        output_pptx_path: Path for output file (default: overwrite input)
    """
    # Convert to absolute paths
    generated_pptx = Path(generated_pptx_path).resolve()
    project_info = Path(project_info_path).resolve()
    
    if not generated_pptx.exists():
        print(f"Error: Generated PPTX not found: {generated_pptx}")
        print(f"Current working directory: {Path.cwd()}")
        return False
    
    if not project_info.exists():
        print(f"Warning: Project info not found: {project_info}. Skipping architecture template slide.")
        deployment_method = None
    else:
        deployment_method = get_deployment_method(project_info)
    
    # Set output path
    if output_pptx_path is None:
        output_pptx_path = generated_pptx
    else:
        output_pptx_path = Path(output_pptx_path)
    
    # Paths to reference files
    script_dir = Path(__file__).parent
    ref_dir = script_dir.parent / "ref"  # ref directory is in template2slide folder
    system_arch_path = ref_dir / "System_architecture.pptx"
    # Try both possible file names
    available_slides_path = ref_dir / "AvailableSlide11.pptx"
    if not available_slides_path.exists():
        available_slides_path = ref_dir / "Available _Slide.pptx"
    
    if not system_arch_path.exists():
        print(f"Warning: System_architecture.pptx not found at {system_arch_path}")
        system_arch_path = None
    
    if not available_slides_path.exists():
        print(f"Warning: Available _Slide.pptx not found at {available_slides_path}")
        available_slides_path = None
    
    # Load presentations
    print(f"Loading generated presentation: {generated_pptx}")
    generated_pres = Presentation(str(generated_pptx))
    
    total_slides_before = len(generated_pres.slides)
    print(f"Original presentation has {total_slides_before} slides")
    
    # Create a copy for output
    if output_pptx_path != generated_pptx:
        shutil.copy2(generated_pptx, output_pptx_path)
        pres = Presentation(str(output_pptx_path))
    else:
        # Create a temporary file to work with
        temp_path = generated_pptx.parent / f"{generated_pptx.stem}_temp.pptx"
        shutil.copy2(generated_pptx, temp_path)
        pres = Presentation(str(temp_path))
    
    # Load reference presentations
    system_arch_pres = None
    available_pres = None
    
    if system_arch_path:
        system_arch_pres = Presentation(str(system_arch_path))
    if available_slides_path:
        available_pres = Presentation(str(available_slides_path))
    
    # Track slides to reorder
    slides_to_reorder = []
    
    # Step 1: Insert architecture template slide after architecture diagram slide
    arch_template_slide_index = None
    if system_arch_path and deployment_method:
        arch_slide_index = DEPLOYMENT_SLIDE_MAP.get(deployment_method)
        
        if arch_slide_index is not None and arch_slide_index < len(system_arch_pres.slides):
            print(f"\nStep 1: Inserting architecture template slide for deployment: {deployment_method}")
            print(f"  Using slide {arch_slide_index + 1} from System_architecture.pptx")
            
            # Find architecture diagram slide in generated presentation
            arch_diagram_index = find_architecture_slide_index(pres)
            print(f"  Architecture diagram slide found at index {arch_diagram_index}")
            
            # Copy slide from system_arch_pres
            new_slide = copy_slide_from_other_pres(system_arch_pres, arch_slide_index, pres)
            arch_template_slide_index = len(pres.slides) - 1
            slides_to_reorder.append((arch_template_slide_index, arch_diagram_index + 1))
            print(f"  ✓ Architecture template slide copied (will be moved to position {arch_diagram_index + 2})")
        else:
            print(f"\nStep 1: Skipping architecture template slide (deployment method '{deployment_method}' not in template)")
    else:
        print("\nStep 1: Skipping architecture template slide (file or deployment method not found)")
    
    # Step 2: Insert Available slides 2-10 after slide 1 (Title slide) - continuously, in order
    # IMPORTANT: These slides must be inserted AFTER slide 1 (Title slide at index 0)
    # Target positions: 1, 2, 3, ..., 9 (0-indexed) which are positions 2, 3, 4, ..., 10 (1-indexed)
    # Slide 1 (Title) is at index 0
    # So Available slides should be at indices 1, 2, 3, ..., 9 (0-indexed)
    # NOTE: We preserve the original backgrounds of reference slides (no copying from generated slides)
    available_slides_2_10_indices = []
    if available_pres and len(available_pres.slides) >= 10:
        print(f"\nStep 2: Inserting Available slides 2-10 after slide 1 (Title slide, continuously)")
        # Insert all slides first, then reorder them to positions 1-9 (0-indexed) = 2-10 (1-indexed)
        # After slide 1 (index 0), so first Available slide is at index 1
        for i in range(1, 10):  # Slides 2-10 from Available (indices 1-9)
            new_slide = copy_slide_from_other_pres(available_pres, i, pres)
            new_index = len(pres.slides) - 1
            # Target position: 1, 2, 3, ..., 9 (0-indexed) = positions 2, 3, 4, ..., 10 (1-indexed)
            # After slide 1 (index 0), so first Available slide is at index 1
            target_position = i  # 1, 2, 3, ..., 9 (0-indexed positions, right after slide 1)
            available_slides_2_10_indices.append((new_index, target_position))
            print(f"  ✓ Copied Available slide {i + 1} (will be moved to position {target_position + 1})")
        print(f"  ✓ All Available slides 2-10 copied (will be inserted continuously after slide 1 - Title slide)")
    else:
        print("\nStep 2: Skipping Available slides 2-10 (file not found or not enough slides)")
    
    # Step 3: Insert Available slides 11-25 at the end
    # NOTE: We preserve the original backgrounds of reference slides (no copying from generated slides)
    if available_pres and len(available_pres.slides) >= 25:
        print(f"\nStep 3: Inserting Available slides 11-25 after last slide")

        for i in range(10, 25):  # Slides 11-25 (indices 10-24)
            new_slide = copy_slide_from_other_pres(available_pres, i, pres)
            print(f"  ✓ Copied Available slide {i + 1}")
        print(f"  ✓ All Available slides 11-25 copied")
    else:
        print("\nStep 3: Skipping Available slides 11-25 (file not found or not enough slides)")
    
    # Reorder slides: move slides to correct positions
    print(f"\nReordering slides...")
    slides = pres.slides._sldIdLst
    
    # Collect all reorder operations
    # IMPORTANT: Available slides 2-10 must be inserted AFTER slide 1 (Title slide, index 0)
    # Architecture template slide will be inserted later and should not affect Available slides position
    all_reorders = []
    
    # For Available slides 2-10: sort by target position (ascending) to ensure correct order
    # This ensures they are inserted continuously: 2, 3, 4, ..., 10 (right after slide 1)
    available_slides_sorted = sorted(available_slides_2_10_indices, key=lambda x: x[1])
    all_reorders.extend(available_slides_sorted)
    
    # Architecture template slide should be inserted after Available slides are in place
    # So we add it after Available slides reorders
    if arch_template_slide_index is not None and len(slides_to_reorder) > 0:
        all_reorders.append((arch_template_slide_index, slides_to_reorder[0][1]))
    
    # IMPORTANT: Sort by target position (ascending), then by source index (ascending)
    # Move from low target positions to high target positions
    # When moving from high source index to low target index, we need to move in order
    # from the first slide (lowest source index) to maintain correct order
    all_reorders.sort(key=lambda x: (x[1], x[0]))
    
    # Get current slide elements as a list
    slide_elements = list(slides)
    
    # Track current positions as we move slides (indices change after each move)
    current_positions = {i: i for i in range(len(slide_elements))}
    
    for source_idx, target_idx in all_reorders:
        if source_idx < len(slide_elements):
            # Get the actual current position of this slide (may have changed due to previous moves)
            actual_source = current_positions.get(source_idx, source_idx)
            if actual_source >= len(slides):
                continue
                
            slide_element = slides[actual_source]
            
            # Remove from current position
            slides.remove(slide_element)
            
            # Insert at target position
            insert_pos = min(target_idx, len(slides))
            slides.insert(insert_pos, slide_element)
            
            # Update positions: all slides that were after the source position move back by 1
            # All slides that are at or after the target position move forward by 1
            for pos in current_positions:
                if current_positions[pos] > actual_source:
                    current_positions[pos] -= 1
                if current_positions[pos] >= insert_pos:
                    current_positions[pos] += 1
            current_positions[source_idx] = insert_pos
            
            print(f"  ✓ Moved slide from position {source_idx + 1} to position {target_idx + 1}")
    
    # Save final presentation
    if output_pptx_path == generated_pptx:
        pres.save(str(generated_pptx))
        if temp_path.exists():
            temp_path.unlink()
    else:
        pres.save(str(output_pptx_path))
    
    total_slides_after = len(pres.slides)
    print(f"\n✓ Complete! Presentation now has {total_slides_after} slides (was {total_slides_before})")
    print(f"  Saved to: {output_pptx_path}")
    
    return True


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python insert_reference_slides.py <generated_pptx> <project_info_json> [output_pptx]")
        print("\nExample:")
        print("  python insert_reference_slides.py output/proposal.pptx output/project_info.json output/final_proposal.pptx")
        sys.exit(1)
    
    generated_pptx = sys.argv[1]
    project_info = sys.argv[2]
    output_pptx = sys.argv[3] if len(sys.argv) > 3 else None
    
    success = insert_reference_slides(generated_pptx, project_info, output_pptx)
    sys.exit(0 if success else 1)

