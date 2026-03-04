#!/usr/bin/env python3
"""
Insert Available slides into generated presentation using python-pptx
Based on pptx skill's proven rearrange.py approach

Usage:
    python insert_available_slides_pro.py <generated.pptx> <AvailableSlide11.pptx> <output.pptx>

Structure:
    1. Slide 1 from generated (Title Slide)
    2. Available Slides 2-10 (9 slides) - inserted after title
    3. Remaining slides from generated (Slides 2 onwards)
    4. Available Slides 11-25 (15 slides) - inserted at end
"""

import sys
import shutil
from pathlib import Path
from copy import deepcopy
import six
from pptx import Presentation


def copy_slide_between_presentations(source_pres, source_idx, dest_pres):
    """
    Copy a slide from source presentation to destination.
    Uses the proven approach from pptx skill's rearrange.py
    """
    source_slide = source_pres.slides[source_idx]

    # Use source's layout to preserve formatting
    new_slide = dest_pres.slides.add_slide(source_slide.slide_layout)

    # Collect all image and media relationships from the source slide
    image_rels = {}
    for rel_id, rel in six.iteritems(source_slide.part.rels):
        if "image" in rel.reltype or "media" in rel.reltype:
            image_rels[rel_id] = rel

    # CRITICAL: Clear placeholder shapes to avoid duplicates
    for shape in new_slide.shapes:
        sp = shape.element
        sp.getparent().remove(sp)

    # Copy all shapes from source
    for shape in source_slide.shapes:
        el = shape.element
        new_el = deepcopy(el)
        new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")

        # Handle picture shapes - need to update the blip reference
        # Look for all blip elements (they can be in pic or other contexts)
        blips = new_el.xpath(".//a:blip[@r:embed]")
        for blip in blips:
            old_rId = blip.get(
                "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
            )
            if old_rId in image_rels:
                # Create a new relationship in the destination slide for this image
                old_rel = image_rels[old_rId]
                # get_or_add returns the rId directly, or adds and returns new rId
                new_rId = new_slide.part.rels.get_or_add(
                    old_rel.reltype, old_rel._target
                )
                # Update the blip's embed reference to use the new relationship ID
                blip.set(
                    "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed",
                    new_rId,
                )

    # Copy any additional image/media relationships that might be referenced elsewhere
    for rel_id, rel in image_rels.items():
        try:
            new_slide.part.rels.get_or_add(rel.reltype, rel._target)
        except Exception:
            pass  # Relationship might already exist

    return new_slide


def insert_available_slides(generated_path, available_path, output_path):
    """
    Insert Available slides into generated presentation.
    """
    print("📊 Inserting Available slides using python-pptx approach...\n")

    # Copy generated presentation as base
    shutil.copy2(generated_path, output_path)
    dest_pres = Presentation(output_path)
    source_pres = Presentation(available_path)

    generated_slides = len(dest_pres.slides)
    available_slides = len(source_pres.slides)

    print(f"  Generated slides: {generated_slides}")
    print(f"  Available slides: {available_slides}\n")

    # Track insertion points
    # After slide 0 (title): insert Available slides 2-10 (indices 1-9)
    # After all generated: insert Available slides 11-25 (indices 10-24)

    insert_count = 0

    # Step 1: Insert Available slides 2-10 after title (slide 0)
    print("📋 Step 1: Inserting Available slides 2-10 after title...")
    for i in range(1, min(10, available_slides)):
        copy_slide_between_presentations(source_pres, i, dest_pres)
        insert_count += 1
        print(f"  ✓ Inserted Available slide {i+1}")

    # Step 2: Insert Available slides 11-25 at the end
    print(f"\n📋 Step 2: Inserting Available slides 11-25 at end...")
    for i in range(10, available_slides):
        copy_slide_between_presentations(source_pres, i, dest_pres)
        insert_count += 1
        print(f"  ✓ Inserted Available slide {i+1}")

    # Save the presentation
    dest_pres.save(output_path)

    total_slides = len(dest_pres.slides)
    print(f"\n✅ Created merged presentation: {output_path}")
    print(f"📊 Total slides: {total_slides}")
    print(f"   - Generated: {generated_slides}")
    print(f"   - Available inserted: {insert_count}")
    print(f"   - Expected: {generated_slides + insert_count}")

    return total_slides


def main():
    if len(sys.argv) < 4:
        print("Usage: python insert_available_slides_pro.py <generated.pptx> <AvailableSlide11.pptx> <output.pptx>")
        print("\nExample:")
        print("  python insert_available_slides_pro.py output_bromma_v2/presentation.pptx \\")
        print("      template2slide-pro/scripts/AvailableSlide11.pptx \\")
        print("      output_bromma_v2/presentation_complete.pptx")
        sys.exit(1)

    generated_path = sys.argv[1]
    available_path = sys.argv[2]
    output_path = sys.argv[3]

    try:
        insert_available_slides(generated_path, available_path, output_path)
        sys.exit(0)
    except Exception as e:
        print(f"\n❌ Error: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()
